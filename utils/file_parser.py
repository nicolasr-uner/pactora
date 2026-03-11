import io
import logging
import threading

_log = logging.getLogger("pactora")

# ─── OCR rate-limiting global ──────────────────────────────────────────────────
# Plan gratuito Gemini: 20 req/día para gemini-2.5-flash.
# Reservamos máx 5 llamadas OCR por sesión de indexación para no agotar la quota.
_OCR_MAX_PER_SESSION = 5
_ocr_calls_remaining = _OCR_MAX_PER_SESSION
_ocr_lock = threading.Lock()


def reset_ocr_quota():
    """Resetea el contador OCR al inicio de cada indexación."""
    global _ocr_calls_remaining
    with _ocr_lock:
        _ocr_calls_remaining = _OCR_MAX_PER_SESSION
    _log.info("[file_parser] OCR quota reseteada — %d llamadas disponibles", _OCR_MAX_PER_SESSION)


def _ocr_acquire() -> bool:
    """Intenta consumir una llamada OCR. Retorna False si se agotó la quota."""
    global _ocr_calls_remaining
    with _ocr_lock:
        if _ocr_calls_remaining <= 0:
            return False
        _ocr_calls_remaining -= 1
        _log.info("[file_parser] OCR call consumida — quedan %d", _ocr_calls_remaining)
        return True


def extract_text_from_file(file_obj, filename: str, gemini_api_key: str = None) -> str:
    """
    Extrae texto de PDF, DOCX, PPTX, XLSX, XLS, CSV, TXT, PNG, JPG, TIFF usando librerías locales.
    Para PDFs escaneados e imágenes intenta Gemini Vision OCR como fallback.
    """
    fname = filename.lower()

    if fname.endswith(".docx"):
        return _extract_docx(file_obj)
    elif fname.endswith(".pdf"):
        file_bytes = file_obj.read() if hasattr(file_obj, "read") else bytes(file_obj)
        return _extract_pdf_bytes(file_bytes)
    elif fname.endswith(".pptx") or fname.endswith(".ppt"):
        return _extract_pptx(file_obj)
    elif fname.endswith(".xlsx") or fname.endswith(".xls"):
        return _extract_excel(file_obj, fname)
    elif fname.endswith(".csv"):
        return _extract_csv(file_obj)
    elif fname.endswith(".txt") or fname.endswith(".md"):
        try:
            raw = file_obj.read() if hasattr(file_obj, "read") else bytes(file_obj)
            return raw.decode("utf-8", errors="replace")
        except Exception:
            return ""
    elif any(fname.endswith(ext) for ext in (".png", ".jpg", ".jpeg", ".tiff", ".tif")):
        file_bytes = file_obj.read() if hasattr(file_obj, "read") else bytes(file_obj)
        return _extract_image_ocr(file_bytes, fname)
    else:
        try:
            raw = file_obj.read() if hasattr(file_obj, "read") else bytes(file_obj)
            return raw.decode("utf-8", errors="replace")
        except Exception:
            return ""


def _extract_docx(file_obj) -> str:
    try:
        from docx import Document
        doc = Document(file_obj)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"Error extrayendo DOCX: {e}"


def _extract_excel(file_obj, fname: str) -> str:
    """Extrae texto de XLSX (openpyxl) o XLS (xlrd)."""
    file_bytes = file_obj.read() if hasattr(file_obj, "read") else bytes(file_obj)

    # XLSX con openpyxl
    if fname.endswith(".xlsx"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Hoja: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None and str(c).strip())
                    if row_text:
                        parts.append(row_text)
            return "\n".join(parts)
        except Exception as e:
            return f"Error extrayendo XLSX: {e}"

    # XLS con xlrd
    try:
        import xlrd
        wb = xlrd.open_workbook(file_contents=file_bytes)
        parts = []
        for sheet in wb.sheets():
            parts.append(f"[Hoja: {sheet.name}]")
            for rx in range(sheet.nrows):
                row_text = " | ".join(str(v) for v in sheet.row_values(rx) if str(v).strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"Error extrayendo XLS: {e}"


def _extract_csv(file_obj) -> str:
    try:
        import csv
        raw = file_obj.read() if hasattr(file_obj, "read") else bytes(file_obj)
        text = raw.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        return "\n".join(" | ".join(row) for row in reader if any(c.strip() for c in row))
    except Exception as e:
        return f"Error extrayendo CSV: {e}"


def _extract_pptx(file_obj) -> str:
    """Extrae texto de presentaciones PPTX (y PPT vía conversión)."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
        file_bytes = file_obj.read() if hasattr(file_obj, "read") else bytes(file_obj)
        prs = Presentation(io.BytesIO(file_bytes))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"[Diapositiva {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except Exception as e:
        return f"Error extrayendo PPTX: {e}"


def _extract_image_ocr(file_bytes: bytes, fname: str) -> str:
    """
    Extrae texto de imágenes (PNG, JPG, TIFF) via Gemini Vision OCR.
    Fallback: retorna cadena vacía si no hay LLM disponible o se agotó la quota.
    """
    try:
        from core.llm_service import LLM_AVAILABLE, GEMINI_API_KEY, _GEMINI_MODEL
        if not LLM_AVAILABLE:
            _log.warning("[file_parser] Imagen recibida pero Gemini no disponible: %s", fname)
            return ""

        if not _ocr_acquire():
            _log.warning("[file_parser] OCR quota agotada — omitiendo imagen: %s", fname)
            return ""

        from google import genai
        from google.genai import types

        if fname.endswith(".png"):
            mime = "image/png"
        elif fname.endswith((".jpg", ".jpeg")):
            mime = "image/jpeg"
        elif fname.endswith((".tiff", ".tif")):
            mime = "image/tiff"
        else:
            mime = "image/png"

        client = genai.Client(api_key=GEMINI_API_KEY)
        response = client.models.generate_content(
            model=_GEMINI_MODEL,
            contents=[
                types.Part(text="Extrae todo el texto visible en esta imagen. Devuelve solo el texto preservando la estructura y párrafos. Si no hay texto, describe brevemente el contenido."),
                types.Part.from_bytes(data=file_bytes, mime_type=mime),
            ],
        )
        result = response.text if response.text else ""
        if result.strip():
            _log.info("[file_parser] OCR imagen OK: %s — %d chars", fname, len(result))
        return result
    except Exception as e:
        err_str = str(e)
        if "429" in err_str or "RESOURCE_EXHAUSTED" in err_str:
            _log.warning("[file_parser] OCR imagen quota 429 — omitiendo: %s", fname)
        else:
            _log.warning("[file_parser] OCR imagen falló para %s: %s", fname, err_str[:120])
        return ""


def _extract_pdf_bytes(file_bytes: bytes) -> str:
    """Intenta pypdf → PyPDF2 → pymupdf (mejor para PDFs complejos)."""
    # Intento 1: pypdf
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        parts = [t for page in reader.pages for t in [page.extract_text()] if t and t.strip()]
        if parts:
            return "\n".join(parts)
    except Exception:
        pass

    # Intento 2: PyPDF2
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        parts = [t for page in reader.pages for t in [page.extract_text()] if t and t.strip()]
        if parts:
            return "\n".join(parts)
    except Exception:
        pass

    # Intento 3: pymupdf (fitz) — mejor extracción en PDFs con layouts complejos
    try:
        import fitz  # pymupdf
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        parts = []
        for page in doc:
            text = page.get_text("text")
            if text and text.strip():
                parts.append(text)
        doc.close()
        if parts:
            return "\n".join(parts)
    except Exception:
        pass

    # Intento 4: Gemini Vision OCR — solo para PDFs escaneados sin texto seleccionable
    try:
        import time as _time
        from core.llm_service import LLM_AVAILABLE, GEMINI_API_KEY, _GEMINI_MODEL
        if not LLM_AVAILABLE:
            return ""

        import fitz  # pymupdf
        from google import genai  # type: ignore
        from google.genai import types  # type: ignore

        client = genai.Client(api_key=GEMINI_API_KEY)
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        parts = []

        for page_num, page in enumerate(doc):
            # Verificar quota antes de cada página
            if not _ocr_acquire():
                _log.warning(
                    "[file_parser] OCR quota agotada en página %d/%d — %d págs procesadas",
                    page_num, len(doc), len(parts)
                )
                break

            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            try:
                response = client.models.generate_content(
                    model=_GEMINI_MODEL,
                    contents=[
                        types.Part(text="Extrae todo el texto de esta página de documento. Devuelve solo el texto preservando estructura y párrafos."),
                        types.Part.from_bytes(data=img_bytes, mime_type="image/png"),
                    ],
                )
                page_text = response.text if response.text else ""
                if page_text.strip():
                    parts.append(page_text)
                # 3s entre páginas para conservar quota diaria
                if page_num < len(doc) - 1:
                    _time.sleep(3.0)
            except Exception as page_err:
                err_str = str(page_err)
                if "429" in err_str or "RESOURCE_EXHAUSTED" in err_str:
                    _log.warning(
                        "[file_parser] OCR 429 en página %d — quota diaria agotada, %d págs recuperadas",
                        page_num, len(parts)
                    )
                    break
                else:
                    _log.warning("[file_parser] OCR página %d falló: %s", page_num, err_str[:80])

        doc.close()

        if parts:
            result = "\n".join(parts)
            _log.info("[file_parser] OCR Gemini Vision exitoso — %d páginas, %d chars", len(parts), len(result))
            return result
    except Exception as e:
        _log.warning("[file_parser] OCR Gemini Vision falló: %s", e)

    return ""
